from pathlib import Path

from file2md.core.base import BaseConverter, ConvertResult
from file2md.utils.markdown import table_to_markdown


class PptxConverter(BaseConverter):
    """
    .pptx / .ppt 파일을 마크다운으로 변환한다.

    변환 규칙:
        - 슬라이드별로 ## Slide N: 제목 섹션을 만든다.
        - 제목 플레이스홀더는 슬라이드 제목으로 사용한다.
        - 본문 텍스트 박스의 내용을 순서대로 추출한다.
        - 테이블이 있으면 마크다운 테이블로 변환한다.
        - 슬라이드 노트(Notes)가 있으면 > 인용 블록으로 추가한다.

    의존성:
        pip install python-pptx
    """

    @property
    def supported_formats(self) -> tuple:
        return ("pptx", "ppt")

    def convert(self, source: str) -> ConvertResult:
        path = Path(source)

        if not path.exists():
            return ConvertResult.failure(source, "pptx", f"파일을 찾을 수 없습니다: {source}")

        try:
            from pptx import Presentation
            from pptx.util import Pt
            from pptx.enum.shapes import MSO_SHAPE_TYPE
        except ImportError:
            return ConvertResult.failure(source, "pptx", "python-pptx가 설치되어 있지 않습니다. pip install python-pptx")

        try:
            prs = Presentation(str(path))
        except Exception as e:
            return ConvertResult.failure(source, "pptx", f"파일을 열 수 없습니다: {e}")

        metadata = {
            "title": path.stem,
            "source": str(path.resolve()),
            "format": "pptx",
            "slides": len(prs.slides),
        }
        frontmatter = self.build_frontmatter(metadata)

        parts: list[str] = []
        total_tables = 0
        total_images = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            # 슬라이드 제목 추출
            slide_title = self._get_title(slide)
            heading = f"## Slide {slide_num}"
            if slide_title:
                heading += f": {slide_title}"
            parts.append(f"\n{heading}\n")

            # 도형 순회
            for shape in slide.shapes:
                shape_type_name = type(shape).__name__

                # 텍스트 박스 / 제목 제외한 일반 텍스트
                if shape.has_text_frame:
                    # 제목 플레이스홀더는 이미 heading에 사용함
                    if self._is_title_placeholder(shape):
                        continue
                    text = self._extract_text_frame(shape.text_frame)
                    if text:
                        parts.append(text)

                # 테이블
                elif shape.has_table:
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        rows.append(cells)
                    md_table = table_to_markdown(rows)
                    if md_table:
                        parts.append(md_table)
                        total_tables += 1

                # 이미지
                elif hasattr(shape, "image"):
                    alt = shape.name or "image"
                    parts.append(f"![{alt}](image)")
                    total_images += 1

            # 슬라이드 노트
            notes = self._get_notes(slide)
            if notes:
                parts.append(f"\n> **Notes:** {notes}\n")

        markdown = frontmatter + self.sanitize("\n".join(parts))

        return ConvertResult(
            markdown_content=markdown,
            source=source,
            format="pptx",
            page_count=len(prs.slides),
            table_count=total_tables,
            image_count=total_images,
        )

    @staticmethod
    def _get_title(slide) -> str:
        """슬라이드에서 제목 텍스트를 반환한다."""
        try:
            title_shape = slide.shapes.title
            if title_shape and title_shape.has_text_frame:
                return title_shape.text.strip()
        except Exception:
            pass
        return ""

    @staticmethod
    def _is_title_placeholder(shape) -> bool:
        """도형이 제목 플레이스홀더인지 확인한다."""
        try:
            from pptx.enum.text import PP_ALIGN
            return shape == shape.part.slide.shapes.title
        except Exception:
            return False

    @staticmethod
    def _extract_text_frame(text_frame) -> str:
        """텍스트 프레임의 모든 단락을 추출한다."""
        lines = []
        for para in text_frame.paragraphs:
            text = para.text.strip()
            if text:
                # 들여쓰기 레벨에 따라 목록 형식 적용
                level = para.level
                prefix = "  " * level + "- " if level > 0 else ""
                lines.append(f"{prefix}{text}")
        return "\n".join(lines)

    @staticmethod
    def _get_notes(slide) -> str:
        """슬라이드 노트를 추출한다."""
        try:
            notes_slide = slide.notes_slide
            if notes_slide:
                text = notes_slide.notes_text_frame.text.strip()
                return text
        except Exception:
            pass
        return ""
